import os
import uuid
import json
from typing import Dict, List, Any, Annotated, Tuple, TypedDict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from langgraph.graph import StateGraph, END
import operator
from langchain_core.messages import AnyMessage, HumanMessage, AIMessage, ToolMessage, SystemMessage
from langchain_openai import AzureChatOpenAI
from langchain_community.tools.tavily_search import TavilySearchResults

# Define the state for the PPT generation graph
class PPTState(TypedDict):
    messages: Annotated[List[AnyMessage], operator.add]
    topic: str
    search_results: List[Dict[str, Any]]
    slides_content: Dict[str, Any]
    ppt_path: str

class LangGraphPPTGenerator:
    def __init__(self, model, search_tool, presentations_dir="presentations"):
        """Initialize the PowerPoint generation tool using LangGraph.
        
        Args:
            model: LLM model to use
            search_tool: Tool to search for information
            presentations_dir: Directory to save generated presentations
        """
        self.model = model
        self.search_tool = search_tool
        self.presentations_dir = presentations_dir
        
        # Create output directory if it doesn't exist
        os.makedirs(self.presentations_dir, exist_ok=True)
        print("Inside ppt class ????????????????????????????????????????????????????????????????")
        
        # Create the graph for PPT generation
        self.graph = self.create_graph()
    
    def create_graph(self):
        """Create the LangGraph for PPT generation."""
        # Define the graph
        graph = StateGraph(PPTState)
        
        # Add nodes
        graph.add_node("search", self.search_for_content)
        graph.add_node("generate_content", self.generate_slides_content)
        graph.add_node("create_ppt", self.create_powerpoint)
        
        # Add edges
        graph.add_edge("search", "generate_content")
        graph.add_edge("generate_content", "create_ppt")
        graph.add_edge("create_ppt", END)
        
        # Set entry point
        graph.set_entry_point("search")
        
        # Compile the graph
        return graph.compile()
    
    async def search_for_content(self, state: PPTState) -> PPTState:
        """Search for content related to the topic."""
        topic = state.get("topic", "")
        
        # Add human message about search
        search_msg = HumanMessage(content=f"Searching for information about '{topic}'...")
        
        # Use the search tool to find information
        search_results = await self.search_tool.ainvoke(topic)
        
        # Add a tool message with the search results
        tool_msg = ToolMessage(
            content=f"Found information about '{topic}'",
            tool_call_id="search-1",
            name="tavily_search_results",
            additional_kwargs={"search_results": search_results}
        )
        
        # Update the state
        return {
            "messages": state.get("messages", []) + [search_msg, tool_msg],
            "topic": topic,
            "search_results": search_results,
            "slides_content": state.get("slides_content", {}),
            "ppt_path": state.get("ppt_path", "")
        }
    
    async def generate_slides_content(self, state: PPTState) -> PPTState:
        """Generate slides content based on search results."""
        topic = state.get("topic", "")
        search_results = state.get("search_results", [])
        
        # Create a prompt for the LLM to generate slide content
        system_prompt = """
        You are an expert presentation designer. Create a comprehensive PowerPoint presentation outline on the provided topic.
        Use the search results to inform your content. Follow these guidelines:

        1. Create a well-structured 10-slide presentation
        2. Include a compelling title slide
        3. Organize content logically with clear slide titles
        4. Keep bullet points concise and meaningful (3-5 per slide)
        5. Include a summary or conclusion slide at the end

        Format your response as a JSON object with this structure:
        {
          "title": "Main Presentation Title",
          "slides": [
            {
              "title": "Slide Title",
              "content": ["Bullet point 1", "Bullet point 2", "..."]
            }
          ]
        }

        The content should be factual, based on the search results, and well-organized.
        DO NOT make up information - stick to what's provided in the search results.
        """
        
        # Format search results into a string
        search_content = ""
        for i, result in enumerate(search_results):
            search_content += f"\n\nSource {i+1}: {result.get('title', '')}\n{result.get('content', '')}"
        
        # Create the human message with the search results
        human_msg = HumanMessage(
            content=f"Create a PowerPoint presentation about '{topic}' based on these search results: {search_content}"
        )
        
        # Get the slides content from the model
        messages = [
            SystemMessage(content=system_prompt),
            human_msg
        ]
        slides_content_msg = await self.model.ainvoke(messages)
        
        # Try to parse the JSON response
        try:
            content_text = slides_content_msg.content
            # Find the JSON part in the response if it's embedded in other text
            json_start = content_text.find('{')
            json_end = content_text.rfind('}') + 1
            
            if json_start >= 0 and json_end > json_start:
                json_str = content_text[json_start:json_end]
                slides_content = json.loads(json_str)
            else:
                # Fallback to a basic template if JSON parsing fails
                slides_content = {
                    "title": f"Presentation on {topic}",
                    "slides": [
                        {"title": "Introduction", "content": [f"Overview of {topic}"]},
                        {"title": "Key Points", "content": ["Important information about the topic"]},
                        {"title": "Conclusion", "content": ["Summary of findings"]}
                    ]
                }
        except Exception as e:
            print(f"Error parsing JSON response: {e}")
            # Create a basic structure if JSON parsing fails
            slides_content = {
                "title": f"Presentation on {topic}",
                "slides": [
                    {"title": "Introduction", "content": [f"Overview of {topic}"]},
                    {"title": "Key Points", "content": ["Important information about the topic"]},
                    {"title": "Conclusion", "content": ["Summary of findings"]}
                ]
            }
        
        # Add AI message with the generated content
        ai_msg = AIMessage(
            content=f"Generated slides content for '{topic}'",
            additional_kwargs={"slides_content": slides_content}
        )
        
        # Update the state
        return {
            "messages": state["messages"] + [human_msg, ai_msg],
            "topic": topic,
            "search_results": search_results,
            "slides_content": slides_content,
            "ppt_path": state.get("ppt_path", "")
        }
    
    async def create_powerpoint(self, state: PPTState) -> PPTState:
        """Create a PowerPoint presentation from the generated content."""
        topic = state.get("topic", "")
        slides_content = state.get("slides_content", {})
        
        # Generate a unique filename
        sanitized_topic = topic.replace(' ', '_').replace('/', '_').replace('\\', '_')[:30]
        filename = f"{sanitized_topic}_{uuid.uuid4().hex[:8]}.pptx"
        file_path = os.path.join(self.presentations_dir, filename)
        
        # Create the presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = slides_content.get("title", f"Presentation on {topic}")
        subtitle_shape.text = "Generated Presentation"
        
        # Style the title slide
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        
        # Add content slides
        for slide_data in slides_content.get("slides", []):
            slide_title = slide_data.get("title", "")
            content_points = slide_data.get("content", [])
            
            layout = prs.slide_layouts[1]  # Title and content (default)
            slide = prs.slides.add_slide(layout)
            
            # Add title
            if hasattr(slide.shapes, "title") and slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = slide_title
                title_shape.text_frame.paragraphs[0].font.size = Pt(36)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
            
            # Add content
            if content_points:
                content_shape = None
                
                # Find the content placeholder
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 1:  # 1 is the content placeholder
                        content_shape = shape
                        break
                
                if content_shape:
                    tf = content_shape.text_frame
                    tf.clear()  # Clear any default text
                    
                    for i, point in enumerate(content_points):
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(52, 73, 94)
                        # First paragraph doesn't need space before
                        if i > 0:
                            p.space_before = Pt(12)
        
        # Add a final "Thank You" slide
        thank_you_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = thank_you_slide.shapes.title
        title_shape.text = "Thank You"
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        
        # Save the presentation
        prs.save(file_path)
        
        # Add a tool message with the PowerPoint file path
        tool_msg = ToolMessage(
            content=f"PowerPoint presentation created and saved at: {file_path}",
            tool_call_id="create-ppt-1",
            name="create_powerpoint",
            additional_kwargs={"file_path": file_path, "filename": filename}
        )
        
        # Update the state
        return {
            "messages": state["messages"] + [tool_msg],
            "topic": topic,
            "search_results": state.get("search_results", []),
            "slides_content": slides_content,
            "ppt_path": file_path
        }
    
    async def generate_presentation(self, topic: str) -> Dict[str, Any]:
        """Generate a PowerPoint presentation on the given topic."""
        # Initialize the state
        initial_state = {
            "messages": [HumanMessage(content=f"Create a PowerPoint presentation about {topic}")],
            "topic": topic,
            "search_results": [],
            "slides_content": {},
            "ppt_path": ""
        }
        
        # Run the graph
        result = await self.graph.ainvoke(initial_state)
        
        # Return the final state
        return result