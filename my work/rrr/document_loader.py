import os
import zipfile
import traceback
import uuid
import pandas as pd
import shutil
from docx import Document as DocxDocument
from pptx import Presentation
from PyPDF2 import PdfReader
import numpy as np
from datetime import datetime
from typing import List


class Document:
    """
    Represents a processed document with content and metadata.
    """
    def __init__(self, page_content: str, metadata: dict):
        self.page_content = page_content
        self.metadata = metadata


class DocumentLoader:
    """
    A versatile document loader for processing various file types.
    """

    def __init__(self):
        """Initialize document loader."""
        self.temp_dirs = []

    def load_documents(self, zip_path: str) -> List[Document]:
        """
        Load documents from a ZIP file. This is a wrapper around load_zip
        to maintain compatibility with the RAGPipeline class.

        Args:
            zip_path (str): Path to the ZIP file.

        Returns:
            List[Document]: List of loaded documents.
        """
        # Simply call the existing load_zip method
        return self.load_zip(zip_path)

    def process_docx(self, file_path: str) -> str:
        """
        Extract text from a Word (DOCX) file.

        Args:
            file_path (str): Path to the DOCX file.

        Returns:
            str: Extracted text content.
        """
        try:
            doc = DocxDocument(file_path)
            text_content = []

            
            core_properties = doc.core_properties
            if core_properties:
                text_content.append("Document Properties:")
                if core_properties.author:
                    text_content.append(f"Author: {core_properties.author}")
                if core_properties.created:
                    text_content.append(f"Created: {core_properties.created}")
                if core_properties.modified:
                    text_content.append(f"Modified: {core_properties.modified}")
                if core_properties.title:
                    text_content.append(f"Title: {core_properties.title}")
                text_content.append("")

            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    
                    if paragraph.style.name and paragraph.style.name != 'Normal':
                        text_content.append(f"[{paragraph.style.name}]")
                    text_content.append(paragraph.text.strip())
                    text_content.append("")  

            
            for table in doc.tables:
                text_content.append("Table Content:")
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    text_content.append(" | ".join(row_text))
                text_content.append("")  

            return "\n".join(text_content)
        except Exception as e:
            print(f"Error processing DOCX file: {e}")
            return ""

    def process_pdf(self, file_path: str) -> str:
        """
        Extract text from a PDF file.

        Args:
            file_path (str): Path to the PDF file.

        Returns:
            str: Extracted text content.
        """
        try:
            reader = PdfReader(file_path)
            text_content = []

            
            metadata = reader.metadata
            if metadata:
                text_content.append("Document Metadata:")
                for key, value in metadata.items():
                    if value and str(value).strip():
                        text_content.append(f"{key}: {value}")
                text_content.append("\n")

            
            for i, page in enumerate(reader.pages, 1):
                text_content.append(f"Page {i}:")
                text = page.extract_text()
                if text.strip():
                    text_content.append(text.strip())
                text_content.append("")  

            return "\n".join(text_content)
        except Exception as e:
            print(f"Error processing PDF file: {e}")
            return ""

    def process_ppt(self, file_path: str) -> str:
        """
        Extract text from a PowerPoint (PPT or PPTX) file.

        Args:
            file_path (str): Path to the PowerPoint file.

        Returns:
            str: Extracted text content.
        """
        try:
            presentation = Presentation(file_path)
            text_content = []

            for slide_number, slide in enumerate(presentation.slides, start=1):
                text_content.append(f"Slide {slide_number}:")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text_content.append(paragraph.text.strip())
                text_content.append("")  

            return "\n".join(text_content)
        except Exception as e:
            print(f"Error processing PowerPoint file: {e}")
            return ""

    def process_csv(self, file_path: str) -> str:
        """
        Process a CSV file and convert it to structured text with enhanced analysis.

        Args:
            file_path (str): Path to the CSV file.

        Returns:
            str: Processed text content.
        """
        try:
            df = pd.read_csv(file_path)
            text_content = []

    
            text_content.append("Dataset Overview:")
            text_content.append(f"Total Records: {len(df)}")
            text_content.append(f"Total Features: {len(df.columns)}")
            text_content.append("")

            
            text_content.append("Column Analysis:")
            for col in df.columns:
                text_content.append(f"\nColumn: {col}")
                dtype = df[col].dtype
                text_content.append(f"Type: {dtype}")

                
                if np.issubdtype(dtype, np.number):
                    stats = df[col].describe()
                    text_content.append("Statistical Summary:")
                    text_content.append(f"Mean: {stats['mean']:.2f}")
                    text_content.append(f"Std Dev: {stats['std']:.2f}")
                    text_content.append(f"Min: {stats['min']:.2f}")
                    text_content.append(f"Max: {stats['max']:.2f}")
                    text_content.append(f"Null Values: {df[col].isnull().sum()}")

                
                elif dtype == 'object':
                    unique_vals = df[col].nunique()
                    text_content.append(f"Unique Values: {unique_vals}")
                    text_content.append(f"Null Values: {df[col].isnull().sum()}")
                    if unique_vals <= 10:  
                        text_content.append("Value Distribution:")
                        for val, count in df[col].value_counts().head().items():
                            text_content.append(f"- {val}: {count}")

                
                elif np.issubdtype(dtype, np.datetime64):
                    text_content.append(f"Date Range: {df[col].min()} to {df[col].max()}")

            
            text_content.append("\nSample Records (First 5):")
            text_content.append(df.head().to_string())

            
            num_cols = df.select_dtypes(include=['int64', 'float64']).columns
            if len(num_cols) > 1:
                text_content.append("\nCorrelation Analysis (Numeric Columns):")
                corr_matrix = df[num_cols].corr()
                for col1 in num_cols:
                    for col2 in num_cols:
                        if col1 < col2:  
                            corr = corr_matrix.loc[col1, col2]
                            if abs(corr) > 0.5:  
                                text_content.append(f"{col1} vs {col2}: {corr:.2f}")

            return "\n".join(text_content)
        except Exception as e:
            print(f"Error processing CSV file: {e}")
            return ""

    def load_zip(self, zip_path: str) -> List[Document]:
        """
        Load documents from a ZIP file with enhanced file support.

        Args:
            zip_path (str): Path to the ZIP file.

        Returns:
            List[Document]: List of loaded documents.
        """
        documents = []

        # Create temp directory for extraction
        extract_path = os.path.join(os.getcwd(), f"extracted_docs_{uuid.uuid4().hex[:8]}")
        self.temp_dirs.append(extract_path)

        try:
            os.makedirs(extract_path, exist_ok=True)

            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                zip_ref.extractall(extract_path)

            for root, _, files in os.walk(extract_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    content = None
                    file_type = file.lower()

                    # Process different file types
                    if file_type.endswith('.csv'):
                        content = self.process_csv(file_path)
                    elif file_type.endswith('.pdf'):
                        content = self.process_pdf(file_path)
                    elif file_type.endswith(('.ppt', '.pptx')):
                        content = self.process_ppt(file_path)
                    elif file_type.endswith('.docx'):
                        content = self.process_docx(file_path)
                    elif file_type.endswith(('.txt', '.md', '.json')):
                        try:
                            with open(file_path, 'r', encoding='utf-8') as f:
                                content = f.read()
                        except UnicodeDecodeError:
                            with open(file_path, 'r', encoding='latin-1') as f:
                                content = f.read()

                    if content and content.strip():
                        # Enhanced metadata
                        metadata = {
                            "source": file_path,
                            "file_type": os.path.splitext(file)[1],
                            "file_name": file,
                            "processed_date": datetime.now().isoformat(),
                            "file_size": os.path.getsize(file_path),
                        }

                        documents.append(Document(
                            page_content=content,
                            metadata=metadata
                        ))

        except Exception as e:
            print(f"Error processing files: {e}")
            traceback.print_exc()

        return documents

    def cleanup(self):
        """Clean up temporary directories."""
        for temp_dir in self.temp_dirs:
            try:
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
            except Exception as e:
                print(f"Error cleaning up directory {temp_dir}: {e}")